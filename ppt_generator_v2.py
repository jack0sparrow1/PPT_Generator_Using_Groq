import os
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from PIL import Image
import io
from dotenv import load_dotenv
import json
import time

# Load environment variables
load_dotenv()

print("✅ All imports successful!")

# ===== Layout Configuration =====
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_LEFT = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
MARGIN_TOP = Inches(0.5)
MARGIN_BOTTOM = Inches(0.5)

CONTENT_TOP = Inches(1.5)   # below title
GUTTER = Inches(0.4)        # space between text and image

IMAGE_MAX_HEIGHT = Inches(4.5)
TEXT_ZONE_RATIO = 0.5 # Percentage of available content width for text
IMAGE_ZONE_RATIO = 0.5 # Percentage of available content width for image


class PPTGenerator:
    def __init__(self):
        """Initialize the PPT Generator with Groq API"""
        self.api_key = os.getenv("GROQ_API_KEY")
        if not self.api_key:
            raise ValueError("GROQ_API_KEY environment variable is required")

        self.client = Groq(api_key=self.api_key)
        self.text_model = "llama-3.3-70b-versatile"
        self.presentation = Presentation()

    def generate_content_outline(self, topic, num_slides=5):
        prompt = f"""Create a professional PowerPoint presentation outline about "{topic}" with EXACTLY {num_slides} slides.

    CRITICAL FORMATTING RULES:
    1. First slide MUST be slide_type: "title" with a title and subtitle
    2. Last slide MUST be slide_type: "conclusion"
    3. Middle slides should be slide_type: "content" or "image_focus"
    4. Each content slide should have 3-5 bullet points
    5. Each bullet point should be a complete sentence (15-20 words)
    6. SEPARATE EACH BULLET POINT WITH ACTUAL NEWLINES (\\n) - NOT commas or periods
    7. Make content informative and detailed, NOT just keywords
    8. Include specific image search terms for visual slides

    Return ONLY a valid JSON array with this EXACT structure:
    [
      {{
        "title": "Main Title Here",
        "subtitle": "Engaging subtitle",
        "content": "",
        "slide_type": "title",
        "image_query": ""
      }},
      {{
        "title": "Content Slide Title",
        "content": "First detailed bullet point as a complete sentence\\nSecond detailed point with explanation\\nThird point with context",
        "slide_type": "content",
        "image_query": "relevant image search term"
      }},
      {{
        "title": "Conclusion",
        "content": "Key takeaway one with details\\nKey takeaway two with context\\nFinal thoughts",
        "slide_type": "conclusion",
        "image_query": ""
      }}
    ]

    IMPORTANT: Use actual newline characters (\\n) to separate bullet points in the "content" field.
    NO markdown, NO code blocks, JUST the JSON array. Make content substantive and informative."""

        try:
            response = self.client.chat.completions.create(
                model=self.text_model,
                messages=[
                    {
                        "role": "system",
                        "content": "You are a presentation expert. Return ONLY valid JSON arrays. Never use markdown formatting. Create detailed, informative content with complete sentences. CRITICAL: Separate bullet points with actual newline characters (\\n), not commas or semicolons."
                    },
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=4000
            )

            raw_content = response.choices[0].message.content.strip()
            
            # Debug: print raw content
            print("Raw response from AI:")
            print(raw_content[:500])
            print("...")
            
            # Remove markdown formatting
            content = raw_content
            if "```json" in content:
                content = content.split("```json")[1].split("```")[0]
            elif "```" in content:
                parts = content.split("```")
                if len(parts) >= 2:
                    content = parts[1]
            
            content = content.strip()
            
            # Debug: print cleaned content
            print("\nCleaned JSON content:")
            print(content[:500])
            print("...")
            
            outline = json.loads(content)
            
            # Validate we got a list
            if not isinstance(outline, list):
                raise ValueError("Response must be a JSON array")
            
            # Ensure first slide is title type
            if outline and outline[0].get("slide_type") != "title":
                outline[0]["slide_type"] = "title"
            
            return outline[:num_slides]  # Ensure we don't exceed requested slides

        except Exception as e:
            print(f"Error generating outline: {e}")
            import traceback
            traceback.print_exc()
            return self._get_fallback_outline(topic, num_slides)
            
    def _get_fallback_outline(self, topic, num_slides):
        return [
            {
                "title": topic,
                "subtitle": "A Comprehensive Overview",
                "content": "",
                "slide_type": "title",
                "image_query": ""
            },
            {
                "title": "Introduction",
                "content": f"This presentation explores the key aspects of {topic}\nWe will examine the fundamental concepts and principles\nUnderstanding these elements is essential for deeper knowledge",
                "slide_type": "content",
                "image_query": f"{topic} introduction"
            },
            {
                "title": "Key Points",
                "content": f"The main characteristics define what makes {topic} important\nHistorical context helps us understand current developments\nModern applications demonstrate practical relevance today",
                "slide_type": "content",
                "image_query": f"{topic} concepts"
            },
            {
                "title": "Applications",
                "content": f"Real-world examples showcase how {topic} is used in practice\nIndustry applications demonstrate tangible benefits and results\nFuture developments point to exciting new possibilities",
                "slide_type": "image_focus",
                "image_query": f"{topic} applications"
            },
            {
                "title": "Conclusion",
                "content": f"We have explored the essential aspects of {topic}\nThese insights provide a foundation for further exploration\nContinued learning will deepen understanding and expertise",
                "slide_type": "conclusion",
                "image_query": ""
            }
        ][:num_slides]

    def download_image(self, query, save_path="temp_image.jpg"):
        try:
            url = "https://api.pexels.com/v1/search"
            headers = {'Authorization': os.getenv('PEXELS_API_KEY')}
            params = {'query': query, 'per_page': 1, 'orientation': 'landscape'}

            response = requests.get(url, headers=headers, params=params, timeout=10)
            response.raise_for_status()

            data = response.json()
            if not data.get('photos'):
                raise ValueError(f"No images found")

            image_url = data['photos'][0]['src']['large']
            img_response = requests.get(image_url, timeout=10)
            img_response.raise_for_status()

            with open(save_path, 'wb') as f:
                f.write(img_response.content)

            return save_path

        except Exception as e:
            # Create placeholder image
            img = Image.new('RGB', (1200, 800), color='#E3F2FD')
            img.save(save_path)
            return save_path

    def _remove_placeholders(self, slide):
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)

    def create_title_slide(self, title, subtitle=""):
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return slide

    def create_content_slide_simple(self, title, content, include_image=False, image_query=None):
        """Even simpler version - uses PowerPoint's default bullet behavior"""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Get content placeholder
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        
        # Clear existing content but KEEP the default formatting (including bullets)
        # Instead of clearing, we'll replace the text in the existing paragraphs
        if "\\n" in content:
            lines = [line.strip() for line in content.split('\\n') if line.strip()]
        else:
            lines = [line.strip() for line in content.split('\n') if line.strip()]
        
        # Clean lines - remove any bullet characters
        cleaned_lines = []
        for line in lines:
            clean_line = line.lstrip('•-*→►▪ ').strip()
            if clean_line:
                cleaned_lines.append(clean_line)
        
        # Set the text using the existing paragraphs to preserve bullet formatting
        for i, paragraph in enumerate(text_frame.paragraphs):
            if i < len(cleaned_lines):
                paragraph.text = cleaned_lines[i]
            else:
                # If we have more paragraphs than lines, clear extra ones
                paragraph.text = ""
        
        # If we have more lines than existing paragraphs, add new ones
        if len(cleaned_lines) > len(text_frame.paragraphs):
            for i in range(len(text_frame.paragraphs), len(cleaned_lines)):
                p = text_frame.add_paragraph()
                p.text = cleaned_lines[i]
        
        # Format all paragraphs
        for paragraph in text_frame.paragraphs:
            if paragraph.text:  # Only format non-empty paragraphs
                paragraph.font.size = Pt(20)
                paragraph.font.name = "Calibri"
                paragraph.space_after = Pt(6)
                paragraph.line_spacing = 1.2
                paragraph.font.color.rgb = RGBColor(51, 51, 51)
                paragraph.level = 0
        
        # Remove any empty paragraphs at the end
        while text_frame.paragraphs and not text_frame.paragraphs[-1].text:
            last_paragraph = text_frame.paragraphs[-1]
            text_frame._element.remove(last_paragraph._element)
        
        # Set position and image (same as before)
        content_shape.left = MARGIN_LEFT
        content_shape.top = CONTENT_TOP
        
        if include_image and image_query:
            try:
                image_path = self.download_image(image_query)
                if image_path and os.path.exists(image_path):
                    usable_width = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
                    text_width = int((usable_width - GUTTER) * TEXT_ZONE_RATIO)
                    image_width = int((usable_width - GUTTER) * IMAGE_ZONE_RATIO)
                    
                    content_shape.width = int(text_width)
                    image_left = int(MARGIN_LEFT + text_width + GUTTER)
                    
                    pic = slide.shapes.add_picture(
                        image_path,
                        image_left,
                        CONTENT_TOP,
                        width=image_width
                    )
                    
                    if pic.height > IMAGE_MAX_HEIGHT:
                        new_height = int(IMAGE_MAX_HEIGHT)
                        new_width = int(pic.width * (new_height / pic.height))
                        pic.height = new_height
                        pic.width = new_width
                    
                    os.remove(image_path)
            except Exception as e:
                print(f"Could not add image: {e}")
        else:
            content_shape.width = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
        
        return slide
    def _get_text_zone_width(self, has_image):
        usable_width = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
        if has_image:
            return (usable_width - GUTTER) * TEXT_ZONE_RATIO
        return usable_width

    def _get_image_zone_width(self):
        usable_width = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
        return (usable_width - GUTTER) * IMAGE_ZONE_RATIO

    def generate_presentation(self, topic, num_slides=5, output_path="presentation.pptx"):
        print(f"Generating {num_slides}-slide presentation on: {topic}")
        
        outline = self.generate_content_outline(topic, num_slides)
        
        # Debug: print the outline to see what content we're getting
        print("\nOutline received:")
        for i, slide_data in enumerate(outline):
            print(f"Slide {i+1}: {slide_data.get('title')}")
            print(f"Content: {slide_data.get('content', '')}")
            print(f"Lines: {len(slide_data.get('content', '').split('\\n'))}")
            print("---")
        
        for i, slide_data in enumerate(outline):
            title = slide_data.get("title", f"Slide {i+1}")
            content = slide_data.get("content", "")
            slide_type = slide_data.get("slide_type", "content")
            subtitle = slide_data.get("subtitle", "")
            image_query = slide_data.get("image_query")

            print(f"\nCreating slide {i+1}: {title} (Type: {slide_type})")
            print(f"Content preview: {content[:50]}...")

            # First slide MUST be title slide
            if i == 0 or slide_type == "title":
                self.create_title_slide(title, subtitle)
            else:
                # For content slides, add images to alternating slides or when specified
                include_image = bool(image_query) or (i % 2 == 1)
                self.create_content_slide_simple(
                    title,
                    content,
                    include_image=include_image,
                    image_query=image_query or title
                )

        self.presentation.save(output_path)
        print(f"\nPresentation saved: {output_path}")
        return output_path

# Initialize the generator
try:
    generator = PPTGenerator()
    print("✅ PPT Generator initialized successfully!")
except ValueError as e:
    print(f"❌ Error: {e}")
    print("Please set your GEMINI_API_KEY first.")

# Generate a presentation
topic = "Why P.Diddy is a good guy?"  # Change this to your desired topic
num_slides = 7  # Change this to your desired number of slides

try:
    output_file = generator.generate_presentation(topic, num_slides, "presentation3.pptx")
except Exception as e:
    print(e)